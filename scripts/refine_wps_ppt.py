from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

BLUE = RGBColor(70, 125, 230)
DARK = RGBColor(17, 24, 39)
GRAY = RGBColor(110, 118, 138)
LIGHT = RGBColor(241, 246, 255)
CYAN = RGBColor(228, 244, 252)
ORANGE = RGBColor(238, 157, 51)
ORANGE_BG = RGBColor(255, 248, 238)
FAINT = RGBColor(186, 208, 250)


def add_box(x, y, w, h, text, fill=None, line=BLUE, font=18, rounded=True, line_w=2.2):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.line.color.rgb = line
    sh.line.width = Pt(line_w)
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    tf = sh.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = 'PingFang SC'
    run.font.size = Pt(font)
    run.font.color.rgb = DARK
    return sh


def add_text(x, y, w, h, text, font=11, color=GRAY, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = 'PingFang SC'
    run.font.size = Pt(font)
    run.font.color.rgb = color
    return tb


def add_line(x1, y1, x2, y2, color=BLUE, width=2.1):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    ln.line.color.rgb = color
    ln.line.width = Pt(width)
    return ln

# tiny caption
add_text(0.35, 0.14, 2.8, 0.25, '根据 /Users/a1-6/Downloads/3.png 精修重绘（形状版）', 9, RGBColor(150, 160, 180))

# top row
add_box(0.88, 0.64, 1.28, 0.42, '租户门户', fill=LIGHT, font=18)
add_box(4.52, 0.64, 1.28, 0.42, '租户门户', fill=LIGHT, font=18)
add_box(8.22, 0.64, 1.28, 0.42, '租户门户', fill=LIGHT, font=18)

# top-right AI generated label
add_box(10.25, 0.18, 1.35, 0.45, '豆包AI生成', fill=ORANGE_BG, line=ORANGE, font=17)

# service box
add_box(9.03, 1.32, 1.95, 0.64, '为租户\n提供人工加AI服务', fill=CYAN, font=15)

# AI data center
add_box(3.98, 1.46, 1.82, 0.52, 'AI数据中心', fill=CYAN, font=20)
add_text(7.52, 1.58, 1.18, 0.2, '平台调用其能力', 10, GRAY)

# mid row
add_box(0.52, 2.88, 1.00, 0.48, '互联网', font=19)
add_box(2.02, 2.98, 1.02, 0.42, 'Agent端', font=16)
add_box(3.86, 2.92, 1.06, 0.45, '胖客户端', font=16)
add_box(5.64, 2.92, 1.06, 0.45, '瘦客户端', font=16)
add_box(7.75, 2.94, 1.06, 0.45, '瘦客户端', font=16)

# bottom blocks
add_box(2.12, 4.74, 0.86, 0.42, '虚拟化', font=16)
add_box(2.20, 5.56, 1.00, 0.42, '网络topo', font=15)
add_box(3.42, 4.74, 0.70, 0.42, 'SDN', font=16)
add_box(4.42, 4.74, 0.70, 0.42, '资产', font=16)
add_box(4.33, 5.55, 0.70, 0.42, '资产', font=16)
add_box(5.78, 4.74, 0.85, 0.42, 'SAAS', font=16)
add_box(7.04, 4.74, 0.70, 0.42, '工单', font=16)
add_box(5.22, 5.25, 1.92, 0.62, '安全模块 | 大数据\n（fw/waf/tw/waflps/scan...）', font=12)
add_box(9.38, 4.96, 1.26, 0.50, 'MSS代运维', font=17)

# vertical guide stems from top portals
add_line(1.52, 1.06, 1.52, 1.34, FAINT, 1.3)
add_line(5.16, 1.06, 5.16, 1.38, FAINT, 1.3)
add_line(8.86, 1.06, 8.86, 1.38, FAINT, 1.3)

# left horizontal chain
add_line(1.52, 3.13, 2.02, 3.18)
add_line(3.04, 3.18, 3.86, 3.15)
add_line(4.92, 3.15, 5.64, 3.15)
add_line(6.70, 3.15, 7.75, 3.16)

# verticals from clients to top
add_line(2.53, 2.98, 2.53, 1.55, FAINT, 1.3)
add_line(4.39, 2.92, 4.39, 2.06, FAINT, 1.3)
add_line(8.28, 2.94, 8.28, 1.96, FAINT, 1.3)
add_line(2.53, 1.55, 1.52, 1.34, FAINT, 1.3)

# AI center to service and callout
add_line(5.80, 1.72, 9.03, 1.62)
add_line(10.00, 1.52, 10.25, 0.63, ORANGE, 1.8)

# AI center down to capability rail
add_line(4.89, 1.98, 4.89, 3.98)
add_line(4.89, 3.98, 7.38, 3.98)
for x in [2.55, 3.77, 4.77, 6.20, 7.39]:
    add_line(x, 3.98, x, 4.74)
add_line(4.68, 3.98, 4.68, 5.55)
add_line(6.18, 3.98, 6.18, 5.25)

# service down to MSS
add_line(10.00, 1.96, 9.95, 4.96)

out = '/Users/a1-6/Downloads/WPS重绘/根据3图重绘_形状版_精修.pptx'
prs.save(out)
print(out)
