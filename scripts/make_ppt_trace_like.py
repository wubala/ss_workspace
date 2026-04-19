from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from PIL import Image, ImageEnhance
from pathlib import Path

orig = Path('/Users/a1-6/Downloads/3.png')
faded = Path('/Users/a1-6/Downloads/WPS重绘/3_底稿淡化.png')
out = Path('/Users/a1-6/Downloads/WPS重绘/根据3图重绘_临摹高仿版.pptx')

# Create faded guide image blended toward white
img = Image.open(orig).convert('RGBA')
white = Image.new('RGBA', img.size, (255,255,255,255))
# keep about 30% original visibility
blend = Image.blend(white, img, 0.34)
blend.save(faded)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])
SW, SH = 13.333, 7.5

slide.shapes.add_picture(str(faded), 0, 0, width=Inches(SW), height=Inches(SH))

BLUE = RGBColor(77,130,235)
DARK = RGBColor(32,47,75)
GRAY = RGBColor(122,132,148)
LIGHT = RGBColor(243,248,255)
CYAN = RGBColor(232,246,252)
ORANGE = RGBColor(241,161,55)
ORANGE_BG = RGBColor(255,248,240)


def rx(x): return x / 2564 * SW
def ry(y): return y / 1440 * SH

def add_box(x, y, w, h, text, fill=None, line=BLUE, font=18, rounded=True, line_w=1.8):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(shape_type, Inches(rx(x)), Inches(ry(y)), Inches(rx(w)), Inches(ry(h)))
    sh.line.color.rgb = line
    sh.line.width = Pt(line_w)
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    tf = sh.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = 'PingFang SC'
    r.font.size = Pt(font)
    r.font.color.rgb = DARK
    return sh

def add_text(x,y,w,h,text,font=10,color=GRAY,align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(rx(x)), Inches(ry(y)), Inches(rx(w)), Inches(ry(h)))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = 'PingFang SC'
    r.font.size = Pt(font)
    r.font.color.rgb = color
    return tb

def add_line(x1,y1,x2,y2,color=BLUE,width=1.6):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(rx(x1)), Inches(ry(y1)), Inches(rx(x2)), Inches(ry(y2)))
    ln.line.color.rgb = color
    ln.line.width = Pt(width)
    return ln

# Overlay editable shapes very close to original positions
add_box(150, 112, 245, 66, '租户门户', fill=LIGHT, font=16)
add_box(835, 112, 245, 66, '租户门户', fill=LIGHT, font=16)
add_box(1495, 110, 250, 68, '租户门户', fill=LIGHT, font=16)
add_box(2205, 52, 300, 70, '豆包AI生成', fill=ORANGE_BG, line=ORANGE, font=17)

add_box(1038, 232, 365, 76, 'AI数据中心', fill=CYAN, font=21)
add_text(1605, 272, 360, 34, '平台调用其能力', 10)
add_box(1990, 225, 340, 104, '为租户\n提供人工加AI服务', fill=CYAN, font=15)

add_box(120, 514, 190, 78, '互联网', font=20)
add_box(565, 545, 195, 68, 'Agent端', font=16)
add_box(927, 532, 202, 68, '胖客户端', font=16)
add_box(1172, 532, 198, 68, '瘦客户端', font=16)
add_box(1585, 534, 198, 68, '瘦客户端', font=16)

add_box(620, 900, 158, 66, '虚拟化', font=16)
add_box(644, 1030, 220, 70, '网络topo', font=15)
add_box(845, 902, 124, 60, 'SDN', font=16)
add_box(1052, 900, 118, 66, '资产', font=16)
add_box(1010, 1028, 118, 66, '资产', font=16)
add_box(1282, 900, 150, 66, 'SAAS', font=16)
add_box(1568, 900, 118, 66, '工单', font=16)
add_box(1238, 986, 426, 122, '安全模块 | 大数据\n（fw/waf/tw/waflps/scan...）', font=12)
add_box(2028, 820, 370, 76, 'MSS代运维', font=18)

# Key lines aligned to original skeleton
for x in [264, 948, 1507]:
    add_line(x, 178, x, 250, RGBColor(184,206,245), 1.0)

add_line(310, 554, 565, 580)
add_line(760, 580, 927, 566)
add_line(1129, 566, 1172, 566)
add_line(1370, 566, 1585, 569)

add_line(330, 540, 330, 246, RGBColor(184,206,245), 1.0)
add_line(330, 246, 264, 250, RGBColor(184,206,245), 1.0)
add_line(1028, 532, 1028, 306, RGBColor(184,206,245), 1.0)
add_line(1684, 534, 1684, 342, RGBColor(184,206,245), 1.0)

add_line(1403, 270, 1990, 275, BLUE, 1.7)
add_line(2330, 225, 2360, 122, ORANGE, 1.5)

add_line(1184, 308, 1184, 758, BLUE, 1.7)
add_line(1184, 758, 1670, 758, BLUE, 1.7)
for x in [700, 907, 1110, 1360, 1627]:
    add_line(x, 758, x, 900, BLUE, 1.6)
add_line(1070, 758, 1070, 1028, BLUE, 1.6)
add_line(1455, 758, 1455, 986, BLUE, 1.6)

add_line(2140, 330, 2210, 820, BLUE, 1.7)

# label
add_text(58, 20, 580, 26, '临摹底稿辅助版：底图可见，便于继续手工微调', 8, RGBColor(155,165,183))

prs.save(out)
print(out)
print(faded)
