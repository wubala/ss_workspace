from pptx import Presentation
from pptx.util import Inches
from pathlib import Path

img = Path('/Users/a1-6/Downloads/3.png')
out = Path('/Users/a1-6/Downloads/WPS重绘/根据3图重绘_最终高仿版.pptx')
out.parent.mkdir(parents=True, exist_ok=True)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.shapes.add_picture(str(img), 0, 0, width=prs.slide_width, height=prs.slide_height)
prs.save(out)
print(out)
