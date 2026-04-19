from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pathlib import Path

out = Path('/Users/a1-6/Downloads/WPS重绘/根据3图重绘_商务版.pptx')
out.parent.mkdir(parents=True, exist_ok=True)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Palette
BG = RGBColor(250, 251, 253)
TITLE = RGBColor(21, 34, 68)
TEXT = RGBColor(51, 65, 85)
MUTED = RGBColor(100, 116, 139)
LINE = RGBColor(203, 213, 225)
BLUE = RGBColor(37, 99, 235)
BLUE_LIGHT = RGBColor(239, 246, 255)
BLUE_LINE = RGBColor(147, 197, 253)
TEAL = RGBColor(14, 116, 144)
TEAL_BG = RGBColor(236, 253, 255)
ORANGE = RGBColor(194, 65, 12)
ORANGE_BG = RGBColor(255, 247, 237)
GREEN = RGBColor(4, 120, 87)
GREEN_BG = RGBColor(236, 253, 245)
GRAY_BG = RGBColor(248, 250, 252)
BORDER = RGBColor(226, 232, 240)


def add_rect(x, y, w, h, fill, line=None, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line
        sh.line.width = Pt(1.2)
    else:
        sh.line.fill.background()
    return sh


def add_box(x, y, w, h, text, fill=BLUE_LIGHT, line=BORDER, font=18, bold=False, color=TEXT):
    sh = add_rect(x, y, w, h, fill, line, True)
    tf = sh.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = 'PingFang SC'
    r.font.size = Pt(font)
    r.font.bold = bold
    r.font.color.rgb = color
    return sh


def add_text(x, y, w, h, text, size=12, color=MUTED, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = 'PingFang SC'
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def add_line(x1, y1, x2, y2, color=BLUE_LINE, width=1.6):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    ln.line.color.rgb = color
    ln.line.width = Pt(width)
    return ln

# Background
bg = add_rect(0, 0, 13.333, 7.5, BG, None, False)
slide.shapes._spTree.remove(bg._element)
slide.shapes._spTree.insert(2, bg._element)

# Header
add_text(0.62, 0.28, 6.0, 0.35, 'AI 数据中心租户服务架构图', 24, TITLE, True)
add_text(0.64, 0.62, 5.5, 0.22, '基于原图业务含义重绘，突出客户接入、能力底座与服务输出关系', 10.5, MUTED)
add_line(0.62, 0.9, 12.7, 0.9, LINE, 0.8)

# Section labels
add_text(0.72, 1.05, 1.8, 0.2, '租户入口', 10.5, BLUE, True)
add_text(0.72, 2.1, 1.8, 0.2, '终端与访问', 10.5, BLUE, True)
add_text(3.95, 2.9, 1.8, 0.2, '能力底座', 10.5, BLUE, True)
add_text(8.85, 1.05, 1.8, 0.2, '服务输出', 10.5, ORANGE, True)
add_text(0.72, 4.4, 1.8, 0.2, '平台能力模块', 10.5, BLUE, True)

# Top entry portals
for x in [0.9, 3.1, 5.3]:
    add_box(x, 1.3, 1.55, 0.52, '租户门户', BLUE_LIGHT, BORDER, 17, True)

# Access row
add_box(0.82, 2.38, 1.08, 0.52, '互联网', GRAY_BG, BORDER, 17, True)
add_box(2.22, 2.38, 1.18, 0.52, 'Agent端', BLUE_LIGHT, BORDER, 17)
add_box(3.82, 2.38, 1.18, 0.52, '胖客户端', BLUE_LIGHT, BORDER, 16)
add_box(5.42, 2.38, 1.18, 0.52, '瘦客户端', BLUE_LIGHT, BORDER, 16)
add_box(7.02, 2.38, 1.18, 0.52, '瘦客户端', BLUE_LIGHT, BORDER, 16)

# Core layer
add_box(3.7, 3.3, 2.0, 0.78, 'AI 数据中心', TEAL_BG, RGBColor(153, 246, 228), 24, True, TEAL)
add_text(4.1, 3.95, 1.2, 0.2, '统一能力底座', 10, TEAL)

# Service output
add_box(9.0, 1.32, 1.95, 0.5, '豆包 AI 生成', ORANGE_BG, RGBColor(253, 186, 116), 17, True, ORANGE)
add_box(8.65, 3.02, 2.35, 0.92, '为租户\n提供人工 + AI 服务', TEAL_BG, RGBColor(153, 246, 228), 19, True, TEAL)
add_text(8.82, 4.0, 1.6, 0.2, '平台调用其能力', 10, MUTED)
add_box(9.1, 5.2, 1.9, 0.64, 'MSS 代运维', GREEN_BG, RGBColor(167, 243, 208), 18, True, GREEN)

# Bottom modules container
panel = add_rect(0.78, 4.7, 7.45, 1.72, RGBColor(255,255,255), BORDER, True)

add_box(1.0, 5.0, 0.96, 0.42, '虚拟化', BLUE_LIGHT, BORDER, 15)
add_box(1.08, 5.66, 1.12, 0.42, '网络 Topo', BLUE_LIGHT, BORDER, 14)
add_box(2.34, 5.0, 0.76, 0.42, 'SDN', BLUE_LIGHT, BORDER, 15)
add_box(3.32, 5.0, 0.76, 0.42, '资产', BLUE_LIGHT, BORDER, 15)
add_box(3.24, 5.66, 0.76, 0.42, '资产', BLUE_LIGHT, BORDER, 15)
add_box(4.42, 5.0, 0.92, 0.42, 'SAAS', BLUE_LIGHT, BORDER, 15)
add_box(5.64, 5.0, 0.76, 0.42, '工单', BLUE_LIGHT, BORDER, 15)
add_box(4.0, 5.62, 2.25, 0.48, '安全模块 / 大数据（fw / waf / tw / waflps / scan ...）', BLUE_LIGHT, BORDER, 11.5)

# Connectors from top portals
for x in [1.68, 3.88, 6.08]:
    add_line(x, 1.82, x, 2.1, LINE, 1.0)

# Access chain
add_line(1.9, 2.64, 2.22, 2.64)
add_line(3.4, 2.64, 3.82, 2.64)
add_line(5.0, 2.64, 5.42, 2.64)
add_line(6.6, 2.64, 7.02, 2.64)

# From access to core
add_line(4.41, 2.9, 4.41, 3.3, LINE, 1.2)
add_line(6.01, 2.9, 6.01, 3.08, LINE, 1.2)
add_line(6.01, 3.08, 5.7, 3.08, LINE, 1.2)

# Core to service
add_line(5.7, 3.68, 8.65, 3.48, BLUE, 1.9)
add_line(9.98, 1.82, 9.98, 3.02, ORANGE, 1.5)

# Core to bottom backbone
add_line(4.7, 4.08, 4.7, 4.82, BLUE, 1.8)
add_line(1.48, 4.82, 6.02, 4.82, BLUE, 1.8)
for x in [1.48, 2.72, 3.7, 4.88, 6.02]:
    add_line(x, 4.82, x, 5.0, BLUE, 1.4)
add_line(3.62, 4.82, 3.62, 5.66, BLUE, 1.4)
add_line(5.14, 4.82, 5.14, 5.62, BLUE, 1.4)

# Service to MSS
add_line(9.98, 3.94, 9.98, 5.2, GREEN, 1.6)

# Right-side note box
note = add_rect(11.35, 1.35, 1.35, 4.5, RGBColor(255,255,255), BORDER, True)
add_text(11.55, 1.58, 0.95, 0.22, '解读要点', 12, TITLE, True)
add_text(11.55, 1.95, 0.95, 3.5,
         '1. 多租户统一通过门户接入\n\n2. Agent / 客户端形成访问入口\n\n3. AI 数据中心承接核心能力\n\n4. 平台向上提供人工 + AI 服务\n\n5. 底部模块为平台基础能力支撑\n\n6. MSS 提供运营与运维输出',
         10, MUTED)

# Footer
add_line(0.62, 6.95, 12.7, 6.95, LINE, 0.8)
add_text(0.65, 7.0, 3.5, 0.18, '商务汇报版 · 适合演示/汇报场景', 8.5, MUTED)

prs.save(out)
print(out)
