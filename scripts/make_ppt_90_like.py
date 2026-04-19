from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

BLUE = RGBColor(77, 130, 235)
DARK = RGBColor(35, 50, 77)
GRAY = RGBColor(125, 135, 150)
LIGHT = RGBColor(243, 248, 255)
CYAN = RGBColor(232, 246, 252)
ORANGE = RGBColor(240, 162, 54)
ORANGE_BG = RGBColor(255, 248, 240)

SW, SH = 13.333, 7.5

def rx(x): return x / 2564 * SW
def ry(y): return y / 1440 * SH

def add_box(x, y, w, h, text, fill=None, line=BLUE, font=18, rounded=True, line_w=2.0):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(shape_type, Inches(rx(x)), Inches(ry(y)), Inches(rx(w)), Inches(ry(h)))
    sh.line.color.rgb = line
    sh.line.width = Pt(line_w)
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    tf = sh.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = 'PingFang SC'
    run.font.size = Pt(font)
    run.font.color.rgb = DARK
    return sh


def add_text(x, y, w, h, text, font=11, color=GRAY, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(rx(x)), Inches(ry(y)), Inches(rx(w)), Inches(ry(h)))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = 'PingFang SC'
    run.font.size = Pt(font)
    run.font.color.rgb = color
    return tb


def add_line(x1, y1, x2, y2, color=BLUE, width=1.8):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(rx(x1)), Inches(ry(y1)), Inches(rx(x2)), Inches(ry(y2)))
    ln.line.color.rgb = color
    ln.line.width = Pt(width)
    return ln

# Small caption
add_text(55, 24, 520, 26, '按原图坐标描摹重建（高相似版）', 8, RGBColor(166,176,192))

# TOP portal row based on OCR positions
add_box(140, 110, 245, 68, '租户门户', fill=LIGHT, font=17)
add_box(820, 110, 245, 68, '租户门户', fill=LIGHT, font=17)
add_box(1370, 108, 255, 72, '租户门户', fill=LIGHT, font=17)

# Top right orange tag and service area
add_box(2198, 52, 295, 72, '豆包AI生成', fill=ORANGE_BG, line=ORANGE, font=18)
add_box(1988, 266, 295, 78, '为租户\n提供人工加AI服务', fill=CYAN, font=15)

# AI data center + platform text
add_box(1020, 245, 365, 78, 'AI数据中心', fill=CYAN, font=21)
add_text(1608, 290, 320, 30, '平台调用其能力', 10, GRAY)

# Mid row boxes
add_box(120, 515, 190, 78, '互联网', font=20)
add_box(585, 550, 190, 64, 'Agent端', font=16)
add_box(950, 535, 196, 64, '胖客户端', font=16)
add_box(1215, 535, 190, 64, '瘦客户端', font=16)
add_box(1600, 540, 190, 64, '瘦客户端', font=16)

# Bottom blocks using OCR-ish positions
add_box(635, 918, 150, 64, '虚拟化', font=16)
add_box(660, 1046, 190, 64, '网络topo', font=15)
add_box(860, 920, 122, 60, 'SDN', font=16)
add_box(1075, 918, 118, 64, '资产', font=16)
add_box(1035, 1040, 118, 64, '资产', font=16)
add_box(1290, 918, 146, 64, 'SAAS', font=16)
add_box(1570, 918, 118, 64, '工单', font=16)
add_box(1238, 1010, 425, 116, '安全模块 | 大数据\n（fw/waf/tw/waflps/scan...）', font=12)
add_box(2042, 842, 360, 76, 'MSS代运维', font=18)

# Vertical guide stems below top portals (faint like original)
for x in [262, 942, 1497]:
    add_line(x, 178, x, 250, RGBColor(173, 198, 244), 1.2)

# Chain lines between mid row boxes
add_line(310, 555, 585, 580)
add_line(775, 580, 950, 567)
add_line(1146, 567, 1215, 567)
add_line(1405, 567, 1600, 570)

# Vertical/diagonal lines from left portal down to Agent chain
add_line(330, 540, 330, 255, RGBColor(173,198,244), 1.2)
add_line(330, 255, 262, 250, RGBColor(173,198,244), 1.2)
# Other client verticals up
add_line(1048, 535, 1048, 328, RGBColor(173,198,244), 1.2)
add_line(1688, 540, 1688, 344, RGBColor(173,198,244), 1.2)

# AI center to service + orange callout
add_line(1385, 284, 1988, 305, BLUE, 1.9)
add_line(2283, 305, 2345, 124, ORANGE, 1.8)

# AI center down to capability backbone
add_line(1180, 323, 1180, 760, BLUE, 1.8)
add_line(1180, 760, 1668, 760, BLUE, 1.8)
for x in [710, 921, 1134, 1363, 1630]:
    add_line(x, 760, x, 918, BLUE, 1.8)
add_line(1094, 760, 1094, 1040, BLUE, 1.8)
add_line(1450, 760, 1450, 1010, BLUE, 1.8)

# service down to MSS
add_line(2142, 344, 2220, 842, BLUE, 1.8)

# a few decorative faint stubs around bottom to mimic original
for x,y1,y2 in [(720,980,1046),(1110,980,1040),(1455,1126,1188),(2230,918,980)]:
    add_line(x,y1,x,y2,RGBColor(173,198,244),1.0)

out = '/Users/a1-6/Downloads/WPS重绘/根据3图重绘_90相似版.pptx'
prs.save(out)
print(out)
