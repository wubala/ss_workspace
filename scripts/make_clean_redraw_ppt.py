from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pathlib import Path

out = Path('/Users/a1-6/Downloads/WPS重绘/根据3图重绘_重新设计版.pptx')
out.parent.mkdir(parents=True, exist_ok=True)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Colors
NAVY = RGBColor(27, 44, 94)
BLUE = RGBColor(64, 114, 230)
MID = RGBColor(100, 149, 237)
LIGHT = RGBColor(239, 246, 255)
CYAN = RGBColor(230, 247, 255)
BORDER = RGBColor(186, 210, 245)
TEXT = RGBColor(31, 41, 55)
SUB = RGBColor(107, 114, 128)
ORANGE = RGBColor(244, 155, 44)
ORANGE_BG = RGBColor(255, 246, 235)
GREEN_BG = RGBColor(236, 253, 245)
GREEN = RGBColor(16, 185, 129)
GRAY_BG = RGBColor(248, 250, 252)


def add_box(x, y, w, h, text, fill=LIGHT, line=BORDER, font=18, bold=False, radius=True, color=TEXT):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line
    sh.line.width = Pt(1.8)
    tf = sh.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = 'PingFang SC'
    r.font.size = Pt(font)
    r.font.bold = bold
    r.font.color.rgb = color
    return sh


def add_text(x, y, w, h, text, size=12, color=SUB, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = 'PingFang SC'
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def add_line(x1, y1, x2, y2, color=BLUE, width=2.0):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    ln.line.color.rgb = color
    ln.line.width = Pt(width)
    return ln

# Background
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
bg.line.fill.background()
slide.shapes._spTree.remove(bg._element)
slide.shapes._spTree.insert(2, bg._element)

# Title
add_text(0.55, 0.22, 5.5, 0.42, 'AI数据中心与租户服务架构（重绘版）', 24, NAVY, True)
add_text(0.58, 0.58, 4.5, 0.24, '按原图业务含义重新整理，适合汇报展示', 11, SUB)

# Top layer title
add_text(0.65, 0.95, 1.8, 0.2, '租户接入层', 11, BLUE, True)
for x in [0.85, 3.05, 5.25]:
    add_box(x, 1.18, 1.55, 0.55, '租户门户', fill=LIGHT, line=BORDER, font=18, bold=True)

# Client access layer
add_text(0.65, 2.0, 1.8, 0.2, '访问与终端层', 11, BLUE, True)
add_box(0.75, 2.28, 1.1, 0.56, '互联网', fill=GRAY_BG, font=18, bold=True)
add_box(2.25, 2.28, 1.2, 0.56, 'Agent端', fill=LIGHT, font=18)
add_box(3.85, 2.28, 1.2, 0.56, '胖客户端', fill=LIGHT, font=17)
add_box(5.45, 2.28, 1.2, 0.56, '瘦客户端', fill=LIGHT, font=17)
add_box(7.05, 2.28, 1.2, 0.56, '瘦客户端', fill=LIGHT, font=17)

# Center core
add_text(3.75, 0.98, 1.8, 0.2, '核心能力层', 11, BLUE, True)
add_box(3.65, 3.25, 2.05, 0.78, 'AI 数据中心', fill=CYAN, line=BLUE, font=24, bold=True)
add_text(4.08, 3.86, 1.2, 0.18, '统一能力底座', 10, BLUE, False, PP_ALIGN.CENTER)

# Right service area
add_text(8.95, 1.0, 1.8, 0.2, '服务输出层', 11, ORANGE, True)
add_box(9.05, 1.15, 1.85, 0.5, '豆包 AI 生成', fill=ORANGE_BG, line=ORANGE, font=18, bold=True, color=RGBColor(146, 64, 14))
add_box(8.65, 2.95, 2.4, 0.88, '为租户\n提供人工 + AI 服务', fill=CYAN, line=BLUE, font=20, bold=True)
add_text(8.78, 3.92, 1.9, 0.18, '平台调用其能力', 10, SUB)
add_box(9.05, 5.18, 1.95, 0.68, 'MSS 代运维', fill=GREEN_BG, line=GREEN, font=19, bold=True, color=RGBColor(6, 95, 70))

# Bottom modules container
add_text(0.65, 4.4, 1.8, 0.2, '平台能力模块', 11, BLUE, True)
container = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.72), Inches(4.62), Inches(7.55), Inches(2.0))
container.fill.solid()
container.fill.fore_color.rgb = RGBColor(250, 252, 255)
container.line.color.rgb = BORDER
container.line.width = Pt(1.2)

add_box(1.0, 4.98, 1.0, 0.46, '虚拟化', fill=LIGHT, font=16)
add_box(1.08, 5.72, 1.18, 0.46, '网络 Topo', fill=LIGHT, font=15)
add_box(2.42, 4.98, 0.82, 0.46, 'SDN', fill=LIGHT, font=16)
add_box(3.45, 4.98, 0.82, 0.46, '资产', fill=LIGHT, font=16)
add_box(3.37, 5.72, 0.82, 0.46, '资产', fill=LIGHT, font=16)
add_box(4.53, 4.98, 0.98, 0.46, 'SAAS', fill=LIGHT, font=16)
add_box(5.83, 4.98, 0.82, 0.46, '工单', fill=LIGHT, font=16)
add_box(4.15, 5.68, 2.25, 0.62, '安全模块 / 大数据\n（fw / waf / tw / waflps / scan ...）', fill=LIGHT, font=13)

# Connectors - top to middle
for start_x, end_x in [(1.62, 2.85), (3.82, 4.45), (6.02, 7.65)]:
    add_line(start_x, 1.73, start_x, 2.06, BORDER, 1.4)
    add_line(start_x, 2.06, end_x, 2.06, BORDER, 1.4)

# Horizontal client chain
add_line(1.85, 2.56, 2.25, 2.56)
add_line(3.45, 2.56, 3.85, 2.56)
add_line(5.05, 2.56, 5.45, 2.56)
add_line(6.65, 2.56, 7.05, 2.56)

# Clients to AI center
add_line(4.45, 2.84, 4.45, 3.25, BORDER, 1.6)
add_line(6.05, 2.84, 6.05, 3.05, BORDER, 1.6)
add_line(6.05, 3.05, 5.7, 3.05, BORDER, 1.6)

# AI to service
add_line(5.7, 3.64, 8.65, 3.39, BLUE, 2.0)
add_line(10.0, 1.65, 10.0, 2.95, ORANGE, 1.8)

# AI to modules backbone
add_line(4.68, 4.03, 4.68, 4.72, BLUE, 1.9)
add_line(1.5, 4.72, 6.25, 4.72, BLUE, 1.9)
for x in [1.5, 2.83, 3.86, 5.02, 6.24]:
    add_line(x, 4.72, x, 4.98, BLUE, 1.5)
add_line(3.78, 4.72, 3.78, 5.72, BLUE, 1.5)
add_line(5.28, 4.72, 5.28, 5.68, BLUE, 1.5)

# Service to MSS
add_line(10.0, 3.83, 10.0, 5.18, GREEN, 1.8)

# Footer note
add_text(8.55, 6.55, 3.7, 0.28, '说明：此页为按原图含义重绘，不追求外观一比一复刻。', 9, SUB)

prs.save(out)
print(out)
